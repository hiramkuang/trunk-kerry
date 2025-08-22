import sys
import os
import pickle
from pathlib import Path
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, 
                             QSplitter, QFileDialog, QPushButton, QComboBox, QTextEdit, 
                             QLabel, QListWidget, QProgressBar, QScrollArea, QGroupBox,
                             QSpinBox, QDoubleSpinBox, QLineEdit, QMessageBox, QCheckBox)
from PyQt5.QtCore import Qt, QSettings, QThread, pyqtSignal
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from unstructured.partition.auto import partition
from langchain.text_splitter import (RecursiveCharacterTextSplitter, 
                                   CharacterTextSplitter)
from sentence_transformers import SentenceTransformer
import numpy as np
import json
import logging
from logging.handlers import RotatingFileHandler
import matplotlib.pyplot as plt
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
import faiss
import spacy
import re
import torch
from transformers import AutoModel, AutoTokenizer

# 配置参数
class Config:
    # 获取当前脚本所在目录
    current_script_path = Path(__file__).resolve()
    
    # 代码所在目录 (RAGSYS/DirectoryLoader)
    CODE_DIR = current_script_path.parent
    
    # 根据目录结构设置路径
    VECTOR_STORE_DIR = CODE_DIR / "vector"  # 向量存储目录
    LOCAL_MODEL_DIR = CODE_DIR / "local_models"  # 本地模型目录
    
    MODEL_NAME = "gpt-4o"
    EMBEDDING_MODEL = "bge-small-zh"
    BACKUP_EMBEDDING_MODEL = "m3e-base"
    CHUNK_SIZE = 1000
    CHUNK_OVERLAP = 200
    OPENROUTER_API_KEY = "your key"
    HF_TOKEN = "hf_..."

    @classmethod
    def init_dirs(cls):
        Path(cls.LOCAL_MODEL_DIR).mkdir(parents=True, exist_ok=True)
        Path(cls.VECTOR_STORE_DIR).mkdir(parents=True, exist_ok=True)

        print(f"模型目录: {cls.LOCAL_MODEL_DIR}")
        print(f"向量库目录: {cls.VECTOR_STORE_DIR}")

# 初始化目录
Config.init_dirs()

# 设置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        RotatingFileHandler('app.log', maxBytes=1024*1024, backupCount=5),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# 自定义嵌入模型类
class CustomEmbeddingModel:
    def __init__(self, model_path):
        self.model = None
        self.tokenizer = None
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        self.load_model(model_path)
    
    def load_model(self, model_path):
        try:
            # 尝试加载为SentenceTransformer模型
            self.model = SentenceTransformer(model_path)
            logger.info(f"加载为SentenceTransformer模型: {model_path}")
        except:
            # 使用transformers加载模型
            try:
                logger.info(f"尝试使用transformers加载模型: {model_path}")
                self.tokenizer = AutoTokenizer.from_pretrained(model_path)
                self.model = AutoModel.from_pretrained(model_path).to(self.device)
                logger.info(f"成功使用transformers加载模型: {model_path}")
            except Exception as e:
                logger.error(f"加载模型失败: {str(e)}")
                raise RuntimeError(f"无法加载模型: {model_path}")
    
    def encode(self, texts):
        if isinstance(self.model, SentenceTransformer):
            return self.model.encode(texts)
        
        # 使用transformers模型生成嵌入
        inputs = self.tokenizer(
            texts, 
            padding=True, 
            truncation=True, 
            return_tensors="pt",
            max_length=512
        ).to(self.device)
        
        with torch.no_grad():
            outputs = self.model(**inputs)
        
        # 使用平均池化获取句子嵌入
        embeddings = self.mean_pooling(outputs, inputs['attention_mask'])
        return embeddings.cpu().numpy()
    
    @staticmethod
    def mean_pooling(model_output, attention_mask):
        token_embeddings = model_output.last_hidden_state
        input_mask_expanded = attention_mask.unsqueeze(-1).expand(token_embeddings.size()).float()
        return torch.sum(token_embeddings * input_mask_expanded, 1) / torch.clamp(input_mask_expanded.sum(1), min=1e-9)

# 后台处理线程
class ProcessingThread(QThread):
    progress = pyqtSignal(int)
    log_message = pyqtSignal(str)
    chunks_ready = pyqtSignal(list)
    error_occurred = pyqtSignal(str)
    
    def __init__(self, files, strategy, strategy_params, embedding_model):
        super().__init__()
        self.files = files
        self.strategy = strategy
        self.strategy_params = strategy_params
        self.embedding_model = embedding_model
        self.chunks = []
    
    def run(self):
        try:
            self.log_message.emit(f"开始处理 {len(self.files)} 个文件")
            self.progress.emit(0)
            
            for i, file_path in enumerate(self.files):
                self.log_message.emit(f"处理文件: {os.path.basename(file_path)}")
                self.progress.emit(int((i / len(self.files)) * 50))
                
                try:
                    text = self.extract_text(file_path)
                    if not text.strip():
                        self.log_message.emit(f"文件 {os.path.basename(file_path)} 内容为空")
                        continue
                        
                    chunks = self.apply_chunking_strategy(text)
                    self.chunks.extend(chunks)
                    
                    self.log_message.emit(f"文件 {os.path.basename(file_path)} 处理完成，生成 {len(chunks)} 个块")
                except Exception as e:
                    self.log_message.emit(f"处理文件 {os.path.basename(file_path)} 时出错: {str(e)}")
                    self.error_occurred.emit(f"处理文件 {os.path.basename(file_path)} 时出错: {str(e)}")
                
                self.progress.emit(int(((i+1) / len(self.files)) * 50))
            
            self.chunks_ready.emit(self.chunks)
            self.progress.emit(100)
            
        except Exception as e:
            self.log_message.emit(f"处理过程中发生错误: {str(e)}")
            self.error_occurred.emit(f"处理过程中发生错误: {str(e)}")
    
    def extract_text(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"不支持的文件类型: {ext}")
            
    def extract_text_from_pdf(self, file_path):
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text
        
    def extract_text_from_docx(self, file_path):
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
        
    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    
    def apply_chunking_strategy(self, text):
        strategy = self.strategy
        
        if strategy == "固定字符数分块":
            chunk_size = self.strategy_params.get('chunk_size', Config.CHUNK_SIZE)
            overlap = self.strategy_params.get('overlap', Config.CHUNK_OVERLAP)
            
            splitter = CharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=overlap
            )
            return splitter.split_text(text)
            
        elif strategy == "递归分块":
            chunk_size = self.strategy_params.get('recursive_size', 500)
            selected_separators = self.strategy_params.get('separators', ["\n\n", "\n", " "])
                
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=int(chunk_size*0.2),
                separators=selected_separators
            )
            return splitter.split_text(text)
            
        elif strategy == "基于格式分块":
            # 使用unstructured库处理结构化文档
            elements = partition(text=text)
            return [str(el) for el in elements]
            
        elif strategy == "基于版式分块":
            max_size = self.strategy_params.get('layout_size', 2000)
            merge_strategy = self.strategy_params.get('merge_strategy', "不合并")
            
            elements = partition(text=text)
            chunks = []
            current_chunk = ""
            
            for el in elements:
                el_text = str(el)
                if len(current_chunk) + len(el_text) <= max_size:
                    current_chunk += "\n" + el_text
                else:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    current_chunk = el_text
                    
            if current_chunk:
                chunks.append(current_chunk.strip())
                
            return chunks
            
        elif strategy == "基于语义分块":
            # 先分句
            if self.strategy_params.get('sentence_model', "spaCy") == "spaCy":
                try:
                    nlp = spacy.load("zh_core_web_sm")
                    doc = nlp(text)
                    sentences = [sent.text for sent in doc.sents]
                except:
                    sentences = [s.strip() for s in re.split(r'(?<=[。！？；…])', text) if s.strip()]
            else:
                sentences = [s.strip() for s in re.split(r'(?<=[。！？；…])', text) if s.strip()]
                
            if not sentences:
                return []
                
            # 计算句子嵌入
            try:
                if self.embedding_model:
                    embeddings = self.embedding_model.encode(sentences)
                else:
                    # 如果没有嵌入模型，使用固定分块
                    splitter = CharacterTextSplitter(
                        chunk_size=500,
                        chunk_overlap=100
                    )
                    return splitter.split_text(text)
            except Exception as e:
                self.log_message.emit(f"计算句子嵌入失败: {str(e)}")
                splitter = CharacterTextSplitter(
                    chunk_size=500,
                    chunk_overlap=100
                )
                return splitter.split_text(text)
            
            # 基于相似度分组
            threshold = self.strategy_params.get('semantic_thresh', 0.85)
            max_group = self.strategy_params.get('max_group', 5)
            
            chunks = []
            current_chunk = []
            
            for i in range(len(sentences)):
                if not current_chunk:
                    current_chunk.append(sentences[i])
                else:
                    # 计算当前句子与最后一个句子的相似度
                    sim = np.dot(embeddings[i], embeddings[i-1])
                    if sim >= threshold and len(current_chunk) < max_group:
                        current_chunk.append(sentences[i])
                    else:
                        chunks.append(" ".join(current_chunk))
                        current_chunk = [sentences[i]]
                        
            if current_chunk:
                chunks.append(" ".join(current_chunk))
                
            return chunks
            
        elif strategy == "命题分块":
            # 使用语言模型识别语义完整的段落
            min_size = self.strategy_params.get('min_topic', 200)
            max_size = self.strategy_params.get('max_topic', 1000)
            
            # 简化的实现
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            chunks = []
            current_chunk = ""
            
            for para in paragraphs:
                if len(current_chunk) + len(para) <= max_size:
                    current_chunk += "\n\n" + para
                else:
                    if current_chunk and len(current_chunk) >= min_size:
                        chunks.append(current_chunk.strip())
                    current_chunk = para
                    
            if current_chunk and len(current_chunk) >= min_size:
                chunks.append(current_chunk.strip())
                
            return chunks
            
        else:
            raise ValueError(f"未知的分块策略: {strategy}")

class VectorStore:
    def __init__(self, index_name="default_index"):
        # 使用os.path.join确保Windows兼容性
        self.index_path = os.path.join(Config.VECTOR_STORE_DIR, f"{index_name}.faiss")
        self.metadata_path = os.path.join(Config.VECTOR_STORE_DIR, f"{index_name}.pkl")
        self.index = None
        self.metadata = []
        
        if os.path.exists(self.index_path):
            self.load_index()
    
    def load_index(self):
        try:
            self.index = faiss.read_index(self.index_path)
            if os.path.exists(self.metadata_path):
                with open(self.metadata_path, 'rb') as f:
                    self.metadata = pickle.load(f)
            logger.info(f"已加载FAISS索引: {self.index_path}")
        except Exception as e:
            logger.error(f"加载FAISS索引失败: {str(e)}")
            self.index = None
            self.metadata = []
    
    def save_index(self):
        try:
            if self.index is not None:
                faiss.write_index(self.index, self.index_path)
                with open(self.metadata_path, 'wb') as f:
                    pickle.dump(self.metadata, f)
                logger.info(f"已保存FAISS索引: {self.index_path}")
        except Exception as e:
            logger.error(f"保存FAISS索引失败: {str(e)}")
    
    def add_vectors(self, vectors, texts):
        if self.index is None:
            dimension = vectors.shape[1]
            self.index = faiss.IndexFlatL2(dimension)
        
        self.index.add(vectors)
        self.metadata.extend(texts)
        self.save_index()
    
    def search(self, query_vector, k=5):
        if self.index is None or self.index.ntotal == 0:
            return []
        
        try:
            distances, indices = self.index.search(query_vector, k)
            results = []
            for i, idx in enumerate(indices[0]):
                if idx >= 0 and idx < len(self.metadata):
                    results.append({
                        'text': self.metadata[idx],
                        'distance': float(distances[0][i])
                    })
            return results
        except Exception as e:
            logger.error(f"向量搜索失败: {str(e)}")
            return []

class ChunkingApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.log_output = []  # 初始化日志输出列表
        self.setWindowTitle("智能文档分块处理工具")
        self.setGeometry(100, 100, 1400, 800)
        
        # 初始化变量
        self.files = []
        self.chunks = []
        self.embeddings = []
        self.current_strategy = None
        self.strategy_params = {}
        self.embedding_model = None
        self.vector_store = VectorStore()
        self.processing_thread = None
        
        # 加载配置
        self.settings = QSettings("RAGSYS", "ChunkingTool")
        self.load_settings()
        
        # 初始化模型
        self.init_models()
        
        # 创建主界面布局
        self.initUI()
        
    def load_settings(self):
        # 从设置加载或使用默认值
        model_path = self.settings.value("model_path", "models", type=str)
        
        # 确保路径是字符串
        if isinstance(model_path, list):
            model_path = model_path[0] if model_path else "models"
        
        # 更新配置
        Config.CHUNK_SIZE = self.settings.value("chunk_size", Config.CHUNK_SIZE, type=int)
        Config.CHUNK_OVERLAP = self.settings.value("chunk_overlap", Config.CHUNK_OVERLAP, type=int)
        Config.EMBEDDING_MODEL = self.settings.value("embedding_model", Config.EMBEDDING_MODEL)
        Config.BACKUP_EMBEDDING_MODEL = self.settings.value("backup_embedding_model", Config.BACKUP_EMBEDDING_MODEL)
        
        # 记录加载的设置
        logger.info(f"加载设置: chunk_size={Config.CHUNK_SIZE}, chunk_overlap={Config.CHUNK_OVERLAP}")
        logger.info(f"嵌入模型: {Config.EMBEDDING_MODEL}, 备用模型: {Config.BACKUP_EMBEDDING_MODEL}")
        logger.info(f"模型路径: {model_path}")
    
    def save_settings(self):
        self.settings.setValue("chunk_size", Config.CHUNK_SIZE)
        self.settings.setValue("chunk_overlap", Config.CHUNK_OVERLAP)
        self.settings.setValue("embedding_model", Config.EMBEDDING_MODEL)
        self.settings.setValue("backup_embedding_model", Config.BACKUP_EMBEDDING_MODEL)
    
    def init_models(self):
        try:
            # 获取模型路径（从设置或默认）
            model_path = self.settings.value("model_path", "models", type=str)
            models_dir = Path(model_path)
            
            # 确保路径是绝对路径
            if not models_dir.is_absolute():
                # 相对于代码目录
                models_dir = Config.CODE_DIR / model_path
            
            logger.info(f"使用模型目录: {models_dir}")
            
            # 验证目录是否存在
            if not models_dir.exists():
                logger.error(f"模型目录不存在: {models_dir}")
                self.log_output.append(f"错误: 模型目录不存在 - {models_dir}")
                return
                
            # 列出目录内容用于调试
            dir_contents = os.listdir(models_dir)
            logger.info(f"模型目录内容: {dir_contents}")
            self.log_output.append(f"模型目录内容: {', '.join(dir_contents)}")
            
            # 尝试加载主模型 - 使用正确的子目录
            main_model_dir = models_dir / Config.EMBEDDING_MODEL
            logger.info(f"尝试加载主模型: {main_model_dir}")
            
            if main_model_dir.exists():
                self.embedding_model = CustomEmbeddingModel(str(main_model_dir))
                logger.info(f"成功加载主模型: {main_model_dir}")
                self.log_output.append(f"已加载模型: {Config.EMBEDDING_MODEL}")
                return
                
            # 尝试加载备用模型 - 使用正确的子目录
            backup_model_dir = models_dir / Config.BACKUP_EMBEDDING_MODEL
            logger.info(f"尝试加载备用模型: {backup_model_dir}")
            
            if backup_model_dir.exists():
                self.embedding_model = CustomEmbeddingModel(str(backup_model_dir))
                logger.info(f"成功加载备用模型: {backup_model_dir}")
                self.log_output.append(f"已加载备用模型: {Config.BACKUP_EMBEDDING_MODEL}")
                return
                
            logger.warning("未找到本地嵌入模型，部分功能可能受限")
            self.log_output.append("警告: 未找到本地嵌入模型，部分功能可能受限")
            
        except Exception as e:
            logger.error(f"无法初始化模型: {str(e)}")
            self.log_output.append(f"模型初始化错误: {str(e)}")
            # 添加堆栈跟踪
            import traceback
            logger.error(traceback.format_exc())

    def initUI(self):
        # 主分割器 - 三栏布局
        main_splitter = QSplitter(Qt.Horizontal)
        
        # 左侧面板 - 文件选择和分块处理
        left_panel = QWidget()
        left_layout = QVBoxLayout()
        
        # 文件选择区
        file_group = QGroupBox("文件选择")
        file_layout = QVBoxLayout()
        
        self.file_list = QListWidget()
        self.btn_select_files = QPushButton("选择文件 (PDF/Word/PPT/TXT)")
        self.btn_select_files.clicked.connect(self.select_files)
        self.btn_clear_files = QPushButton("清空文件列表")
        self.btn_clear_files.clicked.connect(self.clear_files)
        
        file_layout.addWidget(self.btn_select_files)
        file_layout.addWidget(self.btn_clear_files)
        file_layout.addWidget(self.file_list)
        file_group.setLayout(file_layout)
        
        # 分块策略选择区
        strategy_group = QGroupBox("分块策略配置")
        strategy_layout = QVBoxLayout()
        
        self.strategy_combo = QComboBox()
        strategies = [
            "固定字符数分块",
            "递归分块",
            "基于格式分块",
            "基于版式分块",
            "基于语义分块",
            "命题分块"
        ]
        self.strategy_combo.addItems(strategies)
        self.strategy_combo.currentIndexChanged.connect(self.update_strategy_ui)
        
        # 策略参数区 - 动态更新
        self.strategy_params_widget = QWidget()
        self.strategy_params_layout = QVBoxLayout()
        self.strategy_params_widget.setLayout(self.strategy_params_layout)
        
        self.btn_process = QPushButton("执行分块处理")
        self.btn_process.clicked.connect(self.process_files)
        
        strategy_layout.addWidget(QLabel("选择分块策略:"))
        strategy_layout.addWidget(self.strategy_combo)
        strategy_layout.addWidget(self.strategy_params_widget)
        strategy_layout.addWidget(self.btn_process)
        strategy_group.setLayout(strategy_layout)
        
        # 分块可视化区
        viz_group = QGroupBox("分块可视化")
        viz_layout = QVBoxLayout()
        
        self.progress_bar = QProgressBar()
        self.progress_bar.setRange(0, 100)
        
        self.canvas = FigureCanvas(plt.figure())
        
        viz_layout.addWidget(QLabel("处理进度:"))
        viz_layout.addWidget(self.progress_bar)
        viz_layout.addWidget(QLabel("分块分布:"))
        viz_layout.addWidget(self.canvas)
        viz_group.setLayout(viz_layout)
        
        left_layout.addWidget(file_group)
        left_layout.addWidget(strategy_group)
        left_layout.addWidget(viz_group)
        left_panel.setLayout(left_layout)
        
        # 中间面板 - 问答会话区
        center_panel = QWidget()
        center_layout = QVBoxLayout()
        
        chat_group = QGroupBox("文档问答")
        chat_layout = QVBoxLayout()
        
        self.chat_history = QTextEdit()
        self.chat_history.setReadOnly(True)
        
        self.chat_input = QTextEdit()
        self.chat_input.setPlaceholderText("输入您的问题...")
        self.chat_input.setMaximumHeight(100)
        
        self.btn_send = QPushButton("发送")
        self.btn_send.clicked.connect(self.send_message)
        
        chat_layout.addWidget(QLabel("对话历史:"))
        chat_layout.addWidget(self.chat_history)
        chat_layout.addWidget(QLabel("输入问题:"))
        chat_layout.addWidget(self.chat_input)
        chat_layout.addWidget(self.btn_send)
        chat_group.setLayout(chat_layout)
        
        center_layout.addWidget(chat_group)
        center_panel.setLayout(center_layout)
        
        # 右侧面板 - 日志和信息输出
        right_panel = QWidget()
        right_layout = QVBoxLayout()
        
        # 分块信息区
        info_group = QGroupBox("分块信息")
        info_layout = QVBoxLayout()
        
        self.chunk_info = QTextEdit()
        self.chunk_info.setReadOnly(True)
        
        info_layout.addWidget(QLabel("分块统计:"))
        info_layout.addWidget(self.chunk_info)
        info_group.setLayout(info_layout)
        
        # 日志输出区
        log_group = QGroupBox("处理日志")
        log_layout = QVBoxLayout()
        
        self.log_output = QTextEdit()
        self.log_output.setReadOnly(True)
        
        # 重定向日志输出
        class QTextEditLogger(logging.Handler):
            def __init__(self, widget):
                super().__init__()
                self.widget = widget

            def emit(self, record):
                msg = self.format(record)
                self.widget.append(msg)
                
        log_handler = QTextEditLogger(self.log_output)
        log_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        logger.addHandler(log_handler)
        
        log_layout.addWidget(QLabel("实时日志:"))
        log_layout.addWidget(self.log_output)
        log_group.setLayout(log_layout)
        
        right_layout.addWidget(info_group)
        right_layout.addWidget(log_group)
        right_panel.setLayout(right_layout)
        
        # 添加三个主面板到主分割器
        main_splitter.addWidget(left_panel)
        main_splitter.addWidget(center_panel)
        main_splitter.addWidget(right_panel)
        
        # 设置初始比例
        main_splitter.setSizes([400, 600, 400])
        
        self.setCentralWidget(main_splitter)
        
        # 初始化策略参数UI
        self.update_strategy_ui()
        
    def update_strategy_ui(self):
        # 清空当前参数UI
        for i in reversed(range(self.strategy_params_layout.count())): 
            widget = self.strategy_params_layout.itemAt(i).widget()
            if widget is not None:
                widget.setParent(None)
            
        strategy = self.strategy_combo.currentText()
        self.current_strategy = strategy
        
        if strategy == "固定字符数分块":
            self.init_fixed_chunking_ui()
        elif strategy == "递归分块":
            self.init_recursive_chunking_ui()
        elif strategy == "基于格式分块":
            self.init_format_chunking_ui()
        elif strategy == "基于版式分块":
            self.init_layout_chunking_ui()
        elif strategy == "基于语义分块":
            self.init_semantic_chunking_ui()
        elif strategy == "命题分块":
            self.init_topic_chunking_ui()
            
    def init_fixed_chunking_ui(self):
        lbl_chunk_size = QLabel("块大小(字符数):")
        self.spin_chunk_size = QSpinBox()
        self.spin_chunk_size.setRange(100, 5000)
        self.spin_chunk_size.setValue(Config.CHUNK_SIZE)
        
        lbl_overlap = QLabel("重叠大小(字符数):")
        self.spin_overlap = QSpinBox()
        self.spin_overlap.setRange(0, 1000)
        self.spin_overlap.setValue(Config.CHUNK_OVERLAP)
        
        self.strategy_params_layout.addWidget(lbl_chunk_size)
        self.strategy_params_layout.addWidget(self.spin_chunk_size)
        self.strategy_params_layout.addWidget(lbl_overlap)
        self.strategy_params_layout.addWidget(self.spin_overlap)
        
    def init_recursive_chunking_ui(self):
        lbl_chunk_size = QLabel("触发递归的块大小(字符数):")
        self.spin_recursive_size = QSpinBox()
        self.spin_recursive_size.setRange(100, 5000)
        self.spin_recursive_size.setValue(500)
        
        lbl_separators = QLabel("分隔符:")
        self.list_separators = QListWidget()
        
        # 预定义的分隔符
        predefined_separators = {
            "段落": "\n\n",
            "句子": "。",
            "换行": "\n",
            "分号": "; ",
            "逗号": ", ",
            "空格": " ",
            "标题": "\n# ",
            "子标题": "\n## ",
            "列表项": "\n* ",
            "数字列表": "\n1. ",
            "表格": "\n|",
            "代码块": "```"
        }
        
        for desc, sep in predefined_separators.items():
            self.list_separators.addItem(f"{desc}: {json.dumps(sep)}")
            
        self.list_separators.setSelectionMode(QListWidget.MultiSelection)
        
        self.strategy_params_layout.addWidget(lbl_chunk_size)
        self.strategy_params_layout.addWidget(self.spin_recursive_size)
        self.strategy_params_layout.addWidget(lbl_separators)
        self.strategy_params_layout.addWidget(self.list_separators)
        
    def init_format_chunking_ui(self):
        lbl_note = QLabel("此策略将自动检测文档结构元素(标题、段落、列表等)")
        self.strategy_params_layout.addWidget(lbl_note)
        
    def init_layout_chunking_ui(self):
        lbl_chunk_size = QLabel("最大块大小(字符数):")
        self.spin_layout_size = QSpinBox()
        self.spin_layout_size.setRange(100, 5000)
        self.spin_layout_size.setValue(2000)
        
        lbl_merge = QLabel("合并策略:")
        self.combo_merge = QComboBox()
        self.combo_merge.addItems(["不合并", "合并小段落", "合并相似元素"])
        
        self.strategy_params_layout.addWidget(lbl_chunk_size)
        self.strategy_params_layout.addWidget(self.spin_layout_size)
        self.strategy_params_layout.addWidget(lbl_merge)
        self.strategy_params_layout.addWidget(self.combo_merge)
        
    def init_semantic_chunking_ui(self):
        lbl_sentence = QLabel("句子分割模型:")
        self.combo_sentence_model = QComboBox()
        self.combo_sentence_model.addItems(["spaCy", "基本分割"])
        
        lbl_semantic = QLabel("语义相似度阈值:")
        self.spin_semantic_thresh = QDoubleSpinBox()
        self.spin_semantic_thresh.setRange(0.5, 1.0)
        self.spin_semantic_thresh.setValue(0.85)
        self.spin_semantic_thresh.setSingleStep(0.05)
        
        lbl_max_group = QLabel("最大分组句子数:")
        self.spin_max_group = QSpinBox()
        self.spin_max_group.setRange(2, 20)
        self.spin_max_group.setValue(5)
        
        self.strategy_params_layout.addWidget(lbl_sentence)
        self.strategy_params_layout.addWidget(self.combo_sentence_model)
        self.strategy_params_layout.addWidget(lbl_semantic)
        self.strategy_params_layout.addWidget(self.spin_semantic_thresh)
        self.strategy_params_layout.addWidget(lbl_max_group)
        self.strategy_params_layout.addWidget(self.spin_max_group)
        
    def init_topic_chunking_ui(self):
        lbl_min_size = QLabel("最小块大小:")
        self.spin_min_topic = QSpinBox()
        self.spin_min_topic.setRange(50, 1000)
        self.spin_min_topic.setValue(200)
        
        lbl_max_size = QLabel("最大块大小:")
        self.spin_max_topic = QSpinBox()
        self.spin_max_topic.setRange(200, 5000)
        self.spin_max_topic.setValue(1000)
        
        self.strategy_params_layout.addWidget(lbl_min_size)
        self.strategy_params_layout.addWidget(self.spin_min_topic)
        self.strategy_params_layout.addWidget(lbl_max_size)
        self.strategy_params_layout.addWidget(self.spin_max_topic)
        
    def select_files(self):
        options = QFileDialog.Options()
        files, _ = QFileDialog.getOpenFileNames(
            self, "选择文档文件", "", 
            "文档文件 (*.pdf *.docx *.pptx *.txt);;所有文件 (*)", 
            options=options)
            
        if files:
            self.files = files
            self.file_list.clear()
            self.file_list.addItems([os.path.basename(f) for f in files])
            logger.info(f"已选择 {len(files)} 个文件")
            
    def clear_files(self):
        self.files = []
        self.file_list.clear()
        logger.info("已清空文件列表")
        
    def get_strategy_params(self):
        strategy = self.current_strategy
        params = {}
        
        if strategy == "固定字符数分块":
            params = {
                'chunk_size': self.spin_chunk_size.value(),
                'overlap': self.spin_overlap.value()
            }
        elif strategy == "递归分块":
            selected_separators = [json.loads(item.text().split(": ")[1]) 
                                  for item in self.list_separators.selectedItems()]
            
            if not selected_separators:
                selected_separators = ["\n\n", "\n", " "]
                
            params = {
                'recursive_size': self.spin_recursive_size.value(),
                'separators': selected_separators
            }
        elif strategy == "基于版式分块":
            params = {
                'layout_size': self.spin_layout_size.value(),
                'merge_strategy': self.combo_merge.currentText()
            }
        elif strategy == "基于语义分块":
            params = {
                'sentence_model': self.combo_sentence_model.currentText(),
                'semantic_thresh': self.spin_semantic_thresh.value(),
                'max_group': self.spin_max_group.value()
            }
        elif strategy == "命题分块":
            params = {
                'min_topic': self.spin_min_topic.value(),
                'max_topic': self.spin_max_topic.value()
            }
            
        return params
        
    def process_files(self):
        if not self.files:
            QMessageBox.warning(self, "警告", "请先选择文件")
            return
            
        # 获取策略参数
        self.strategy_params = self.get_strategy_params()
        
        # 禁用处理按钮
        self.btn_process.setEnabled(False)
        self.btn_process.setText("处理中...")
        
        # 清空之前的结果
        self.chunks = []
        self.embeddings = []
        self.vector_store = VectorStore()
        
        # 启动后台处理线程
        self.processing_thread = ProcessingThread(
            self.files, 
            self.current_strategy,
            self.strategy_params,
            self.embedding_model
        )
        
        # 连接信号
        self.processing_thread.progress.connect(self.progress_bar.setValue)
        self.processing_thread.log_message.connect(self.log_output.append)
        self.processing_thread.chunks_ready.connect(self.on_chunks_ready)
        self.processing_thread.error_occurred.connect(self.on_processing_error)
        
        # 启动线程
        self.processing_thread.start()
        
    def on_processing_error(self, error):
        QMessageBox.critical(self, "错误", error)
        self.btn_process.setEnabled(True)
        self.btn_process.setText("执行分块处理")
        
    def on_chunks_ready(self, chunks):
        self.chunks = chunks
        self.btn_process.setEnabled(True)
        self.btn_process.setText("执行分块处理")
        
        # 向量化处理
        self.vectorize_chunks()
        
        # 可视化分块结果
        self.visualize_chunks()
        
    def vectorize_chunks(self):
        if not self.chunks:
            QMessageBox.warning(self, "警告", "没有可向量化的文本块")
            return
            
        try:
            if self.embedding_model:
                self.embeddings = self.embedding_model.encode(self.chunks)
                self.vector_store.add_vectors(np.array(self.embeddings), self.chunks)
                self.log_output.append(f"已完成 {len(self.embeddings)} 个文本块的向量化")
            else:
                self.log_output.append("警告: 没有可用的嵌入模型，跳过向量化")
        except Exception as e:
            self.log_output.append(f"向量化过程中出错: {str(e)}")
            QMessageBox.critical(self, "错误", f"向量化过程中出错: {str(e)}")
    
    def visualize_chunks(self):
        if not self.chunks:
            return
            
        # 清除之前的图形
        self.canvas.figure.clear()
        ax = self.canvas.figure.add_subplot(111)
        
        # 计算每个块的长度
        chunk_lengths = [len(chunk) for chunk in self.chunks]
        
        # 绘制块长度分布
        ax.hist(chunk_lengths, bins=20, alpha=0.7, color='skyblue')
        ax.set_xlabel('块长度(字符数)', fontsize=10)
        ax.set_ylabel('数量', fontsize=10)
        ax.set_title('分块长度分布', fontsize=12)
        
        # 添加统计信息
        stats = (
            f"总块数: {len(self.chunks)}\n"
            f"平均长度: {int(np.mean(chunk_lengths))}\n"
            f"最小长度: {min(chunk_lengths)}\n"
            f"最大长度: {max(chunk_lengths)}"
        )
        ax.text(0.95, 0.95, stats, transform=ax.transAxes,
                verticalalignment='top', horizontalalignment='right',
                bbox=dict(facecolor='white', alpha=0.7), fontsize=9)
        
        self.canvas.draw()
        
        # 更新分块信息
        self.update_chunk_info()
        
    def update_chunk_info(self):
        if not self.chunks:
            self.chunk_info.clear()
            return
            
        chunk_lengths = [len(chunk) for chunk in self.chunks]
        info_text = (
            f"分块策略: {self.current_strategy}\n"
            f"总块数: {len(self.chunks)}\n"
            f"平均长度: {int(np.mean(chunk_lengths))} 字符\n"
            f"最小长度: {min(chunk_lengths)} 字符\n"
            f"最大长度: {max(chunk_lengths)} 字符\n\n"
            "前3个块预览:\n"
        )
        
        for i, chunk in enumerate(self.chunks[:3]):
            info_text += f"\n--- 块 {i+1} (长度: {len(chunk)} 字符) ---\n"
            info_text += chunk[:200] + ("..." if len(chunk) > 200 else "") + "\n"
            
        self.chunk_info.setPlainText(info_text)
        
    def send_message(self):
        message = self.chat_input.toPlainText().strip()
        if not message:
            return
            
        self.chat_input.clear()
        self.chat_history.append(f"<b>用户:</b> {message}")
        
        if not self.embedding_model or not self.vector_store or not self.embeddings:
            response = "<b>系统:</b> 错误: 嵌入模型未初始化或向量存储不可用"
        else:
            try:
                # 向量化查询
                query_embedding = self.embedding_model.encode([message])
                
                # 在向量存储中搜索
                results = self.vector_store.search(query_embedding, k=3)
                
                if results:
                    response = "<b>系统:</b> 找到的相关文本块:\n\n"
                    for i, result in enumerate(results):
                        similarity = 1 - result['distance']
                        response += f"<b>结果 {i+1}</b> (相似度: {similarity:.2f}):\n"
                        response += result['text'][:300] + ("..." if len(result['text']) > 300 else "") + "\n\n"
                else:
                    response = "<b>系统:</b> 未找到相关结果"
            except Exception as e:
                response = f"<b>系统:</b> 查询过程中出错: {str(e)}"
            
        self.chat_history.append(response)

    def closeEvent(self, event):
        self.save_settings()
        super().closeEvent(event)

if __name__ == "__main__":
    app = QApplication(sys.argv)
    app.setStyle('Fusion')
    window = ChunkingApp()
    window.show()
    sys.exit(app.exec_())